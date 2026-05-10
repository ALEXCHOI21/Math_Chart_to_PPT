from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
import math
import matplotlib.pyplot as plt
import numpy as np
from scipy.interpolate import make_interp_spline

class PPTConverter:
    def __init__(self):
        self.width = Inches(10)
        self.height = Inches(5.625)
        self.center_x = self.width / 2
        self.center_y = self.height / 2
        self.scale = Inches(0.5)  # 1 unit = 0.5 inch

    def create_ppt(self, data, output_path):
        if isinstance(data, list) and len(data) > 0:
            data = data[0]
        prs = Presentation()
        # 16:9 비율 설정
        prs.slide_width = self.width
        prs.slide_height = self.height
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 배경색 (다크 모드)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0) # Pitch Black to match site

        # 1. 축 그리기 (화살표 포함)
        self._draw_axes(slide, data.get("axes", {}))

        # 2. 곡선 그리기 (Smooth)
        for curve in data.get("curves", []):
            self._draw_curve(slide, curve)

        # 3. 점선 그리기
        for dl in data.get("dashed_lines", []):
            self._draw_dashed_line(slide, dl)

        # 4. 점 그리기
        for point in data.get("points", []):
            self._draw_point(slide, point)

        # 5. 화살표 그리기 (곡선 화살표 지원)
        for arrow in data.get("arrows", []):
            self._draw_arrow(slide, arrow)

        # 6. 라벨 그리기
        for label in data.get("labels", []):
            self._draw_label(slide, label)

        prs.save(output_path)
        return output_path

    def _to_ppt_coords(self, x, y):
        """그래프 좌표를 PPT 슬라이드 좌표로 변환"""
        ppt_x = self.center_x + (x * self.scale)
        ppt_y = self.center_y - (y * self.scale)  # Y축은 위가 양수
        return ppt_x, ppt_y

    def _draw_axes(self, slide, axes):
        x_range = axes.get("x_range", [-5, 5])
        y_range = axes.get("y_range", [-5, 5])
        
        # X축 (Professional thin line)
        sx, sy = self._to_ppt_coords(x_range[0], 0)
        ex, ey = self._to_ppt_coords(x_range[1], 0)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
        line.line.color.rgb = RGBColor(255, 255, 255)
        line.line.width = Pt(1.5)
        line.line.end_arrowhead = 1  # Standard arrow
        line.line.end_arrow_width = 3 # Large
        line.line.end_arrow_length = 3 # Large

        # Y축 (Professional thin line)
        sx, sy = self._to_ppt_coords(0, y_range[0])
        ex, ey = self._to_ppt_coords(0, y_range[1])
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
        line.line.color.rgb = RGBColor(255, 255, 255)
        line.line.width = Pt(1.5)
        line.line.end_arrowhead = 1  # Standard arrow
        line.line.end_arrow_width = 3 # Large
        line.line.end_arrow_length = 3 # Large

    def _draw_curve(self, slide, curve):
        pts = curve.get("points", [])
        if len(pts) < 2:
            return

        # ── High-Fidelity Parametric Spline Interpolation ────────────────
        try:
            import numpy as np
            from scipy.interpolate import make_interp_spline
            pts_arr = np.array(pts)
            
            if len(pts_arr) >= 3:
                # Parametric interpolation: handles non-monotonic curves correctly
                t = np.linspace(0, 1, len(pts_arr))
                t_new = np.linspace(0, 1, 300)
                k = min(3, len(pts_arr) - 1)
                spl_x = make_interp_spline(t, pts_arr[:, 0], k=k)(t_new)
                spl_y = make_interp_spline(t, pts_arr[:, 1], k=k)(t_new)
                draw_pts = list(zip(spl_x, spl_y))
            else:
                draw_pts = pts
        except Exception as e:
            draw_pts = pts

        # FreeformBuilder로 매끄러운 선 생성
        start_x, start_y = self._to_ppt_coords(draw_pts[0][0], draw_pts[0][1])
        builder = slide.shapes.build_freeform(start_x, start_y)
        
        for i in range(1, len(draw_pts)):
            px, py = self._to_ppt_coords(draw_pts[i][0], draw_pts[i][1])
            builder.add_line_segments([(px, py)])
        
        shape = builder.convert_to_shape()
        # 색상 처리
        hex_color = curve.get("color", "#FFFFFF").replace("#", "")
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        
        shape.line.color.rgb = RGBColor(r, g, b)
        shape.line.width = Pt(curve.get("width", 2.5))
        if curve.get("style") == "dashed":
            shape.line.dash_style = 2

    def _draw_dashed_line(self, slide, dl):
        dl_type = dl.get("type", "")
        val = dl.get("val") or dl.get("y", 0) if dl_type == "horizontal" else dl.get("val") or dl.get("x", 0)
        
        if dl_type == "horizontal":
            sx, sy = self._to_ppt_coords(dl.get("start", 0), val)
            ex, ey = self._to_ppt_coords(dl.get("end", 5), val)
        else: # vertical
            sx, sy = self._to_ppt_coords(val, dl.get("start", 0))
            ex, ey = self._to_ppt_coords(val, dl.get("end", 5))
            
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
        line.line.color.rgb = RGBColor(120, 120, 120)
        line.line.dash_style = 2
        line.line.width = Pt(1)

    def _draw_arrow(self, slide, arrow):
        is_curved = arrow.get("is_curved", False)
        pts = arrow.get("points", [])
        
        if is_curved and len(pts) >= 3:
            # 곡선 화살표도 촘촘한 보간 적용
            try:
                import numpy as np
                from scipy.interpolate import make_interp_spline
                pts_arr = np.array(pts)
                t = np.linspace(0, 1, len(pts_arr))
                t_new = np.linspace(0, 1, 100)
                spl_x = make_interp_spline(t, pts_arr[:, 0], k=min(3, len(pts)-1))(t_new)
                spl_y = make_interp_spline(t, pts_arr[:, 1], k=min(3, len(pts)-1))(t_new)
                draw_pts = list(zip(spl_x, spl_y))
            except:
                draw_pts = pts

            sx, sy = self._to_ppt_coords(draw_pts[0][0], draw_pts[0][1])
            builder = slide.shapes.build_freeform(sx, sy)
            for i in range(1, len(draw_pts)):
                px, py = self._to_ppt_coords(draw_pts[i][0], draw_pts[i][1])
                builder.add_line_segments([(px, py)])
            shape = builder.convert_to_shape()
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(1.5)
            shape.line.end_arrowhead = 1
            shape.line.end_arrow_width = 2
            shape.line.end_arrow_length = 2
        else:
            fx, fy = arrow.get("from", [0, 0])
            tx, ty = arrow.get("to", [0, 0])
            sx, sy = self._to_ppt_coords(fx, fy)
            ex, ey = self._to_ppt_coords(tx, ty)
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
            line.line.color.rgb = RGBColor(255, 255, 255)
            line.line.width = Pt(1.2)
            line.line.end_arrowhead = 1
            line.line.end_arrow_width = 2
            line.line.end_arrow_length = 2

    def _draw_point(self, slide, point):
        coord = point.get("coord") or [point.get("x", 0), point.get("y", 0)]
        x, y = self._to_ppt_coords(coord[0], coord[1])
        size = Pt(12) # Increased size
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - size/2, y - size/2, size, size)
        shp.line.color.rgb = RGBColor(255, 255, 255)
        shp.line.width = Pt(1.0)
        
        if point.get("type") == "hollow":
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(0, 0, 0) # Background match
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(255, 255, 255)

    def _draw_label(self, slide, label):
        pos = label.get("pos") or [label.get("x", 0), label.get("y", 0)]
        x, y = self._to_ppt_coords(pos[0], pos[1])
        
        # 텍스트 상자 추가 (자동 크기 조정 지원)
        tx = slide.shapes.add_textbox(x - Pt(20), y - Pt(25), Pt(120), Pt(40))
        tf = tx.text_frame
        tf.word_wrap = False
        text = label.get("text", "")
        tf.text = text
        for para in tf.paragraphs:
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.font.size = Pt(18) # Increased size
            para.font.bold = True
            para.font.name = "Inter"

    def export_as_png(self, data, output_path):
        """Matplotlib을 사용하여 고해상도 PNG 그래프 생성 (웹 전시용)"""
        if isinstance(data, list) and len(data) > 0:
            data = data[0]

        plt.figure(figsize=(10, 8), facecolor='black')
        ax = plt.gca()
        ax.set_facecolor('black')

        # 축 설정
        axes = data.get("axes", {})
        xr = axes.get("x_range", [-5, 5])
        yr = axes.get("y_range", [-5, 5])
        
        plt.xlim(xr[0] - 0.5, xr[1] + 0.5)
        plt.ylim(yr[0] - 0.5, yr[1] + 0.5)

        # 축 선 (Professional thin white lines)
        plt.axhline(0, color='white', linewidth=1.5)
        plt.axvline(0, color='white', linewidth=1.5)

        # 축 화살표
        plt.arrow(xr[0]-0.2, 0, xr[1]-xr[0]+0.4, 0, head_width=0.15, head_length=0.2, fc='white', ec='white', length_includes_head=True)
        plt.arrow(0, yr[0]-0.2, 0, yr[1]-yr[0]+0.4, head_width=0.15, head_length=0.2, fc='white', ec='white', length_includes_head=True)

        # 축 라벨
        plt.text(xr[1] + 0.3, 0.2, axes.get("x_label", "x"), color='white', fontsize=18, fontweight='bold', fontname='Inter')
        plt.text(0.2, yr[1] + 0.3, axes.get("y_label", "y"), color='white', fontsize=18, fontweight='bold', fontname='Inter')
        plt.text(-0.3, -0.4, axes.get("origin_label", "O"), color='white', fontsize=18, fontweight='bold', fontname='Inter')

        # 곡선 그리기 (Smooth)
        for curve in data.get("curves", []):
            pts = np.array(curve.get("points", []))
            if len(pts) >= 3:
                t = np.linspace(0, 1, len(pts))
                t_new = np.linspace(0, 1, 300)
                spl_x = make_interp_spline(t, pts[:, 0], k=min(3, len(pts)-1))(t_new)
                spl_y = make_interp_spline(t, pts[:, 1], k=min(3, len(pts)-1))(t_new)
                
                color = curve.get("color", "#FFFFFF")
                width = curve.get("width", 2.5)
                ls = '--' if curve.get("style") == "dashed" else '-'
                plt.plot(spl_x, spl_y, color=color, linewidth=width, linestyle=ls)
            elif len(pts) == 2:
                plt.plot(pts[:, 0], pts[:, 1], color=curve.get("color", "#FFFFFF"), linewidth=curve.get("width", 2.5))

        # 점선
        for dl in data.get("dashed_lines", []):
            val = dl.get("val") or (dl.get("y", 0) if dl.get("type") == "horizontal" else dl.get("x", 0))
            if dl.get("type") == "horizontal":
                plt.hlines(val, dl.get("start", 0), dl.get("end", 5), colors='gray', linestyles='dashed', linewidth=1)
            else:
                plt.vlines(val, dl.get("start", 0), dl.get("end", 5), colors='gray', linestyles='dashed', linewidth=1)

        # 점
        for pt in data.get("points", []):
            coord = pt.get("coord") or [pt.get("x", 0), pt.get("y", 0)]
            m = 'o'
            fc = 'black' if pt.get("type") == "hollow" else 'white'
            plt.plot(coord[0], coord[1], marker=m, markerfacecolor=fc, markeredgecolor='white', markersize=10)

        # 화살표
        for arr in data.get("arrows", []):
            if arr.get("is_curved") and "points" in arr:
                pts = np.array(arr["points"])
                plt.plot(pts[:, 0], pts[:, 1], color='#B4B4B4', linewidth=1.5)
                # 마지막 점에 화살표 머리
                plt.annotate('', xy=(pts[-1][0], pts[-1][1]), xytext=(pts[-2][0], pts[-2][1]),
                             arrowprops=dict(arrowstyle='->', color='#B4B4B4', lw=1.5))
            else:
                f = arr.get("from", [0, 0])
                t = arr.get("to", [0, 0])
                plt.annotate('', xy=(t[0], t[1]), xytext=(f[0], f[1]),
                             arrowprops=dict(arrowstyle='->', color='#B4B4B4', lw=1.5))

        # 라벨
        for lbl in data.get("labels", []):
            p = lbl.get("pos") or [lbl.get("x", 0), lbl.get("y", 0)]
            plt.text(p[0], p[1], lbl.get("text", ""), color='white', fontsize=16, fontname='Inter')

        plt.axis('off')
        plt.tight_layout()
        plt.savefig(output_path, dpi=150, bbox_inches='tight', pad_inches=0.1)
        plt.close()
